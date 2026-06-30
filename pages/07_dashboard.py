import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from utils.dashboard_manager import DashboardManager
from pathlib import Path

def show():
    st.markdown('<h1 class="main-header">📐 Customizable Dashboard</h1>', unsafe_allow_html=True)
    
    if st.session_state.data is None:
        st.warning("⚠️ Please upload data first!")
        return
    
    df = st.session_state.data
    manager = DashboardManager()
    
    # Initialize session state
    if 'dashboard_widgets' not in st.session_state:
        st.session_state.dashboard_widgets = []
    if 'dashboard_layout' not in st.session_state:
        st.session_state.dashboard_layout = manager.get_default_layout()
    
    # Tab layout
    tab1, tab2, tab3 = st.tabs(["🎨 Build Dashboard", "📊 View Dashboard", "💾 Saved Dashboards"])
    
    with tab1:
        build_dashboard(df, manager)
    
    with tab2:
        view_dashboard(df)
    
    with tab3:
        saved_dashboards(manager)

def build_dashboard(df, manager):
    st.subheader("🎨 Build Your Dashboard")
    
    col1, col2 = st.columns([1, 2])
    
    with col1:
        st.markdown("### 🧩 Widget Library")
        widget_type = st.selectbox("Widget Type", ["Chart", "Metric", "Table", "Text"])
        
        if widget_type == "Chart":
            chart_type = st.selectbox("Chart Type", ["Bar", "Line", "Scatter", "Histogram", "Box"])
            
            numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
            cat_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
            
            x_col = st.selectbox("X-Axis", cat_cols if cat_cols else numeric_cols)
            y_col = st.selectbox("Y-Axis", numeric_cols)
            title = st.text_input("Chart Title", f"{y_col} by {x_col}")
            
            if st.button("➕ Add Chart Widget"):
                widget = {
                    'type': 'chart',
                    'chart_type': chart_type,
                    'x_col': x_col,
                    'y_col': y_col,
                    'title': title
                }
                st.session_state.dashboard_widgets.append(widget)
                st.success("✅ Chart added!")
                st.rerun()
        
        elif widget_type == "Metric":
            numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
            metric_col = st.selectbox("Select Metric", numeric_cols)
            metric_name = st.text_input("Display Name", metric_col)
            agg_func = st.selectbox("Aggregation", ["Sum", "Mean", "Max", "Min", "Count"])
            
            if st.button("➕ Add Metric Widget"):
                widget = {
                    'type': 'metric',
                    'col': metric_col,
                    'name': metric_name,
                    'agg': agg_func
                }
                st.session_state.dashboard_widgets.append(widget)
                st.success("✅ Metric added!")
                st.rerun()
        
        elif widget_type == "Table":
            cols = st.multiselect("Select Columns", df.columns.tolist(), default=df.columns[:3].tolist())
            rows = st.slider("Rows to Show", 5, 50, 10)
            
            if st.button("➕ Add Table Widget"):
                widget = {
                    'type': 'table',
                    'cols': cols,
                    'rows': rows
                }
                st.session_state.dashboard_widgets.append(widget)
                st.success("✅ Table added!")
                st.rerun()
        
        elif widget_type == "Text":
            text_type = st.selectbox("Text Type", ["Heading", "Paragraph", "Markdown"])
            text_content = st.text_area("Content", "Enter your text here...")
            
            if st.button("➕ Add Text Widget"):
                widget = {
                    'type': 'text',
                    'text_type': text_type,
                    'content': text_content
                }
                st.session_state.dashboard_widgets.append(widget)
                st.success("✅ Text added!")
                st.rerun()
        
        st.markdown("---")
        
        # Layout controls
        st.markdown("### 📐 Layout Settings")
        columns = st.slider("Columns", 1, 4, 3)
        theme = st.selectbox("Theme", ["Light", "Dark"])
        
        if st.button("💾 Save Dashboard Layout"):
            st.session_state.dashboard_layout['columns'] = columns
            st.session_state.dashboard_layout['theme'] = theme
            st.success("✅ Layout saved!")
    
    with col2:
        st.markdown("### 📊 Dashboard Preview")
        
        if st.session_state.dashboard_widgets:
            # Show active widgets
            for idx, widget in enumerate(st.session_state.dashboard_widgets):
                with st.container():
                    col1, col2 = st.columns([6, 1])
                    with col1:
                        display_widget(df, widget)
                    with col2:
                        if st.button("🗑️", key=f"del_{idx}"):
                            st.session_state.dashboard_widgets.pop(idx)
                            st.rerun()
        else:
            st.info("Add widgets from the left panel to build your dashboard")
        
        st.markdown("---")
        
        # Export options
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("📤 Export to PowerPoint", use_container_width=True):
                export_to_pptx(df, st.session_state.dashboard_widgets)
        
        with col2:
            if st.button("🔗 Shareable Link", use_container_width=True):
                create_shareable_link()
        
        with col3:
            if st.button("💾 Save Dashboard", use_container_width=True):
                save_current_dashboard(manager)

def display_widget(df, widget):
    if widget['type'] == 'chart':
        # Create chart based on type
        if widget['chart_type'] == 'Bar':
            fig = px.bar(df, x=widget['x_col'], y=widget['y_col'], title=widget['title'])
        elif widget['chart_type'] == 'Line':
            fig = px.line(df, x=widget['x_col'], y=widget['y_col'], title=widget['title'])
        elif widget['chart_type'] == 'Scatter':
            fig = px.scatter(df, x=widget['x_col'], y=widget['y_col'], title=widget['title'])
        elif widget['chart_type'] == 'Histogram':
            fig = px.histogram(df, x=widget['x_col'], title=widget['title'])
        elif widget['chart_type'] == 'Box':
            fig = px.box(df, x=widget['x_col'], y=widget['y_col'], title=widget['title'])
        
        fig.update_layout(height=300)
        st.plotly_chart(fig, use_container_width=True)
    
    elif widget['type'] == 'metric':
        # Calculate metric
        col = widget['col']
        agg = widget['agg']
        
        if agg == 'Sum':
            value = df[col].sum()
        elif agg == 'Mean':
            value = df[col].mean()
        elif agg == 'Max':
            value = df[col].max()
        elif agg == 'Min':
            value = df[col].min()
        elif agg == 'Count':
            value = len(df)
        
        st.metric(widget['name'], f"{value:,.0f}")
    
    elif widget['type'] == 'table':
        st.dataframe(df[widget['cols']].head(widget['rows']), use_container_width=True)
    
    elif widget['type'] == 'text':
        if widget['text_type'] == 'Heading':
            st.markdown(f"### {widget['content']}")
        elif widget['text_type'] == 'Paragraph':
            st.write(widget['content'])
        elif widget['text_type'] == 'Markdown':
            st.markdown(widget['content'])

def view_dashboard(df):
    st.subheader("📊 Dashboard View")
    
    if st.session_state.dashboard_widgets:
        # Apply layout
        columns = st.session_state.dashboard_layout.get('columns', 3)
        cols = st.columns(columns)
        
        for idx, widget in enumerate(st.session_state.dashboard_widgets):
            with cols[idx % columns]:
                display_widget(df, widget)
    else:
        st.info("No widgets added yet. Go to 'Build Dashboard' tab to add widgets.")

def saved_dashboards(manager):
    st.subheader("💾 Saved Dashboards")
    
    if manager.dashboards:
        for name, dashboard in manager.dashboards.items():
            with st.expander(f"📊 {name}"):
                st.caption(f"Created: {dashboard.get('created_at', 'Unknown')}")
                st.caption(f"Widgets: {len(dashboard.get('widgets', []))}")
                st.caption(f"Layout: {dashboard.get('layout', {}).get('columns', 3)} columns")
                
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("🔄 Load", key=f"load_{name}"):
                        st.session_state.dashboard_widgets = dashboard.get('widgets', [])
                        st.session_state.dashboard_layout = dashboard.get('layout', {})
                        st.success(f"✅ Loaded {name}")
                        st.rerun()
                with col2:
                    if st.button("🗑️ Delete", key=f"del_{name}"):
                        manager.delete_dashboard(name)
                        st.success(f"✅ Deleted {name}")
                        st.rerun()
    else:
        st.info("No saved dashboards. Build and save one above!")

def export_to_pptx(df, widgets):
    try:
        from pptx import Presentation
        from pptx.util import Inches
        import io
        
        prs = Presentation()
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Data Dashboard Report"
        subtitle.text = f"Generated: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Add slides for each widget
        for widget in widgets:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = widget.get('title', 'Dashboard Widget')
            
            # Add content (simplified for demo)
            body = slide.placeholders[1]
            body.text = f"Widget Type: {widget.get('type', 'Unknown')}"
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        st.download_button(
            label="📥 Download PowerPoint",
            data=pptx_bytes,
            file_name="dashboard_report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
        st.success("✅ PowerPoint exported successfully!")
        
    except ImportError:
        st.warning("⚠️ python-pptx not installed. Install with: pip install python-pptx")

def create_shareable_link():
    # Generate shareable link with current filters
    import base64
    import json
    
    # Encode dashboard state
    state = {
        'widgets': st.session_state.dashboard_widgets,
        'layout': st.session_state.dashboard_layout
    }
    
    state_json = json.dumps(state)
    state_b64 = base64.b64encode(state_json.encode()).decode()
    
    share_url = f"{st.get_option('server.baseUrlPath')}?dashboard={state_b64}"
    
    st.info(f"🔗 Shareable Link: {share_url}")
    st.code(share_url)

def save_current_dashboard(manager):
    name = st.text_input("Dashboard Name", "My Dashboard")
    
    if name and st.button("💾 Confirm Save"):
        manager.save_dashboard(
            name,
            st.session_state.dashboard_layout,
            st.session_state.dashboard_widgets,
            st.session_state.dashboard_layout.get('theme', 'light')
        )
        st.success(f"✅ Dashboard '{name}' saved!")
        st.rerun()
